#!/usr/bin/env python3
"""Build the 50+TechBridge three-lesson teaching deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsmap
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# Brand — matches site navy + LMT orange
NAVY = RGBColor(0x08, 0x36, 0x55)
NAVY_DEEP = RGBColor(0x05, 0x24, 0x3A)
INDIGO = RGBColor(0x2B, 0x31, 0x5F)
ORANGE = RGBColor(0xE2, 0x75, 0x2E)
CYAN = RGBColor(0x00, 0xA3, 0xDC)
GREEN = RGBColor(0x65, 0xBD, 0x53)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF7, 0xF4, 0xEF)
SLATE = RGBColor(0x3D, 0x4A, 0x5C)
MUTED = RGBColor(0x5B, 0x6B, 0x7C)

W = Inches(13.333)
H = Inches(7.5)

OUT = Path(__file__).with_name("50plus-techbridge-free-lessons.pptx")


def _set_run(run, text, size, bold=False, color=WHITE, font="Arial"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, l, t, w, h, text, size=28, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p.add_run() if p.runs else p.runs[0] if False else p.add_run(), text, size, bold, color)
    # python-pptx: first paragraph has no run until we add one; clean empty default
    if tf.paragraphs[0].runs and tf.paragraphs[0].runs[0].text == "" and len(tf.paragraphs[0].runs) > 1:
        r0 = tf.paragraphs[0].runs[0]
        r0.text = ""
    return box


def textbox(slide, l, t, w, h, lines, default_size=28, default_bold=False, default_color=WHITE, align=PP_ALIGN.LEFT):
    """lines: list of str or (text, size, bold, color)."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, bold, color = item, default_size, default_bold, default_color
        else:
            text, size, bold, color = item[0], item[1], item[2] if len(item) > 2 else default_bold, item[3] if len(item) > 3 else default_color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(8)
        run = p.add_run()
        _set_run(run, text, size, bold, color)
    return box


def rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def bar(slide, fill=ORANGE):
    rect(slide, Inches(0), Inches(0), W, Inches(0.12), fill)


def footer(slide, page, total, light=False):
    color = RGBColor(0x9A, 0xB0, 0xC0) if not light else MUTED
    textbox(
        slide,
        Inches(0.5),
        Inches(7.12),
        Inches(9.5),
        Inches(0.3),
        [("50+TechBridge  ·  Free Lessons  ·  Always free", 12, False, color)],
    )
    textbox(
        slide,
        Inches(10.4),
        Inches(7.12),
        Inches(2.4),
        Inches(0.3),
        [(f"{page}  /  {total}", 12, False, color)],
        align=PP_ALIGN.RIGHT,
    )


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def new_dark(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, NAVY)
    return s


def new_light(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, CREAM)
    bar(s)
    return s


def card(slide, l, t, w, h, fill, title, body, title_color=WHITE, body_color=None):
    rect(slide, l, t, w, h, fill)
    if body_color is None:
        body_color = RGBColor(0xD6, 0xE6, 0xF0) if fill != CREAM else SLATE
    textbox(slide, l + Inches(0.22), t + Inches(0.18), w - Inches(0.4), Inches(0.45), [(title, 18, True, title_color)])
    textbox(slide, l + Inches(0.22), t + Inches(0.62), w - Inches(0.4), h - Inches(0.8), [(body, 16, False, body_color)])


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    total = 22

    # 1 Title
    s = new_dark(prs)
    rect(s, 0, 0, Inches(0.22), H, ORANGE)
    textbox(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4), [("FREE LESSONS  ·  1 OF 3 SERIES", 16, True, ORANGE)])
    textbox(s, Inches(0.7), Inches(2.0), Inches(12), Inches(2.2), [
        ("You're not bad at technology.", 40, True, WHITE),
        ("You've just never had the right teacher.", 32, False, RGBColor(0xC8, 0xDF, 0xF0)),
    ])
    textbox(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.2), [
        ("Welcome  →  Talk to AI  →  Don't get scammed", 22, True, CYAN),
        ("50+TechBridge  ·  Learn More Technologies  ·  Always free", 16, False, RGBColor(0x9A, 0xB0, 0xC0)),
    ])
    notes(s, "Open with the house line. Pause. Do not rush. You are talking to a capable adult.")

    # 2 Map
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.9), [
        ("The map", 36, True, NAVY),
        ("Three lessons. One job each. Belong  →  Do  →  Protect.", 18, False, SLATE),
    ])
    cards = [
        (ORANGE, "1  Welcome", "You belong here.\nYou are not too late.\nWrite one real question."),
        (CYAN, "2  Talk to AI", "One tool.\nOne real question.\nOne useful answer."),
        (NAVY, "3  Don't get scammed", "Spot a fake.\nName the pressure.\nLock passwords and 2FA."),
    ]
    for i, (color, title, body) in enumerate(cards):
        x = Inches(0.6) + i * Inches(4.15)
        rect(s, x, Inches(1.7), Inches(3.95), Inches(4.5), color)
        textbox(s, x + Inches(0.25), Inches(1.95), Inches(3.45), Inches(1.0), [(title, 22, True, WHITE)])
        textbox(s, x + Inches(0.25), Inches(3.1), Inches(3.45), Inches(2.6), [(body, 20, False, WHITE)])
    footer(s, 2, total, light=True)
    notes(s, "Hold this map. Lesson 2 is the win. Lesson 3 is the lock. Do not add a fourth lesson today.")

    # 3 Who this is for
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.8), [("This is for you if…", 36, True, NAVY)])
    items = [
        "The phone keeps changing and you are tired of feeling behind.",
        "You want to stay useful at work — or find work — and the ads all say AI.",
        "You want to see a grandchild's face and not miss the button.",
        "You are done being the person who 'should have known' about a scam text.",
        "Someone you love sent you the link and said, just try it.",
    ]
    textbox(s, Inches(0.7), Inches(1.4), Inches(12), Inches(5.2), [(f"•  {t}", 22, False, SLATE) for t in items])
    footer(s, 3, total, light=True)
    notes(s, "All of those are good reasons. No experience needed. One device is enough.")

    # 4 Not starting from zero
    s = new_dark(prs)
    bar(s)
    textbox(s, Inches(0.6), Inches(0.4), Inches(12), Inches(1.0), [
        ("You are not starting from zero", 34, True, WHITE),
        ("AARP research, adults 50+  ·  Say the source on screen", 16, False, RGBColor(0x9A, 0xB0, 0xC0)),
    ])
    stats = [
        ("90%", "own a smartphone"),
        ("30%", "have used generative AI (up from 18% last year)"),
        ("8 in 10", "have used some form of AI already"),
    ]
    for i, (n, label) in enumerate(stats):
        x = Inches(0.6) + i * Inches(4.15)
        rect(s, x, Inches(2.0), Inches(3.95), Inches(3.6), INDIGO)
        textbox(s, x + Inches(0.2), Inches(2.3), Inches(3.55), Inches(1.4), [(n, 48, True, ORANGE)], align=PP_ALIGN.CENTER)
        textbox(s, x + Inches(0.25), Inches(3.9), Inches(3.45), Inches(1.4), [(label, 18, False, WHITE)], align=PP_ALIGN.CENTER)
    footer(s, 4, total)
    notes(s, "Two numbers is enough if you are short on time: 90% smartphone, 18% to 30% generative AI. Beginner is allowed.")

    # 5 Lesson 1 title
    s = new_dark(prs)
    rect(s, 0, 0, Inches(0.22), H, ORANGE)
    textbox(s, Inches(0.7), Inches(2.3), Inches(12), Inches(2.4), [
        ("LESSON 1", 18, True, ORANGE),
        ("Welcome", 54, True, WHITE),
        ("You belong here. You are not too late.", 24, False, RGBColor(0xC8, 0xDF, 0xF0)),
    ])
    footer(s, 5, total)
    notes(s, "Read the long Welcome script. This slide is the chapter card.")

    # 6 Fridge 1
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.5), [("FRIDGE LINE", 14, True, ORANGE)])
    textbox(s, Inches(0.6), Inches(2.2), Inches(12), Inches(2.4), [
        ("You are not too late.", 40, True, NAVY),
        ("You just needed the right teacher.", 32, False, SLATE),
    ], align=PP_ALIGN.CENTER)
    footer(s, 6, total, light=True)
    notes(s, "Hold five seconds. This is the photograph slide.")

    # 7 Homework L1
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.9), [
        ("Do this before Lesson 2", 34, True, NAVY),
        ("Write one real question from your life this week.", 18, False, SLATE),
    ])
    examples = [
        "I have chicken, rice, and spinach. Give me three easy dinners.",
        "Help me write a polite email to my HOA about the parking.",
        "Explain this doctor's note in plain English. I covered the name.",
        "I am applying for a job. Give me ten interview questions.",
    ]
    textbox(s, Inches(0.7), Inches(1.5), Inches(12), Inches(4.2), [(f"•  {t}", 20, False, SLATE) for t in examples])
    textbox(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.6), [("Sticky note. Photograph it. That is your ticket.", 18, True, ORANGE)])
    footer(s, 7, total, light=True)
    notes(s, "If their mind goes blank: 'I feel behind. Give me three ways an adult my age can use AI this week without sharing private information.'")

    # 8 Lesson 2 title
    s = new_dark(prs)
    rect(s, 0, 0, Inches(0.22), H, CYAN)
    textbox(s, Inches(0.7), Inches(2.3), Inches(12), Inches(2.6), [
        ("LESSON 2", 18, True, CYAN),
        ("Talk to AI", 54, True, WHITE),
        ("One tool. One real question. One useful answer.", 24, False, RGBColor(0xC8, 0xDF, 0xF0)),
    ])
    footer(s, 8, total)
    notes(s, "Show the real screen after this card. ChatGPT Free or Copilot if it is already on the PC.")

    # 9 Safety
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.9), [
        ("Before anyone types", 34, True, NAVY),
        ("Put this on screen. Do not skip it.", 18, False, SLATE),
    ])
    bans = [
        ("Social Security number", "Never"),
        ("Bank, card, Medicare numbers", "Never"),
        ("Passwords or phone codes", "Never"),
        ("Someone else's private business", "Never"),
    ]
    for i, (t, tag) in enumerate(bans):
        y = Inches(1.5) + i * Inches(1.15)
        rect(s, Inches(0.6), y, Inches(12.1), Inches(1.05), WHITE)
        rect(s, Inches(0.6), y, Inches(0.16), Inches(1.05), ORANGE)
        textbox(s, Inches(1.0), y + Inches(0.22), Inches(9.5), Inches(0.6), [(t, 24, True, NAVY)])
        textbox(s, Inches(10.6), y + Inches(0.28), Inches(1.8), Inches(0.5), [(tag, 18, True, ORANGE)])
    footer(s, 9, total, light=True)
    notes(s, "Health, money, and the law are drafts. Verify with a qualified human.")

    # 10 Pattern
    s = new_dark(prs)
    bar(s, CYAN)
    textbox(s, Inches(0.6), Inches(0.4), Inches(12), Inches(1.0), [
        ("The only pattern you need", 32, True, WHITE),
        ("Situation  +  Ask  +  Format", 28, True, CYAN),
    ])
    textbox(s, Inches(0.7), Inches(2.0), Inches(12), Inches(3.6), [
        ('"I am 68. I cook for two. I have leftover roast chicken, rice, and spinach. Give me three easy dinners. No spicy food. For each dinner, list extra ingredients and the steps in short sentences."', 22, False, WHITE),
    ])
    textbox(s, Inches(0.7), Inches(5.8), Inches(12), Inches(0.8), [
        ("Weak:  \"dinners\"        Strong:  who you are, what you have, what you want back.", 18, False, RGBColor(0xC8, 0xDF, 0xF0)),
    ])
    footer(s, 10, total)
    notes(s, "Type this live at 150% zoom. Then follow up: Make the second dinner even simpler. One pan.")

    # 11 Useful + follow-up
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7), [("A useful answer — and how to steer", 32, True, NAVY)])
    rect(s, Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.0), NAVY)
    textbox(s, Inches(0.85), Inches(1.55), Inches(5.5), Inches(4.5), [
        ("A useful answer", 20, True, ORANGE),
        ("I understand it.", 22, False, WHITE),
        ("I could do the next step.", 22, False, WHITE),
        ("Nothing private was required.", 22, False, WHITE),
        ("If it sounds too sure, check the number or the 'always.'", 18, False, RGBColor(0xC8, 0xDF, 0xF0)),
    ])
    rect(s, Inches(6.85), Inches(1.3), Inches(5.85), Inches(5.0), INDIGO)
    textbox(s, Inches(7.1), Inches(1.55), Inches(5.4), Inches(4.5), [
        ("Follow-ups that work", 20, True, CYAN),
        ("Make that simpler.", 22, False, WHITE),
        ("Give me three options.", 22, False, WHITE),
        ("I have 15 minutes. Cut this down.", 22, False, WHITE),
        ("Explain the second step like I have never done it.", 20, False, WHITE),
    ])
    footer(s, 11, total, light=True)
    notes(s, "The follow-up is the skill. You are not being rude. You are steering.")

    # 12 Lesson 3 title
    s = new_dark(prs)
    rect(s, 0, 0, Inches(0.22), H, ORANGE)
    textbox(s, Inches(0.7), Inches(2.3), Inches(12), Inches(2.6), [
        ("LESSON 3", 18, True, ORANGE),
        ("Don't get scammed", 50, True, WHITE),
        ("Spot a fake. Name the pressure. Lock the door.", 22, False, RGBColor(0xC8, 0xDF, 0xF0)),
    ])
    footer(s, 12, total)
    notes(s, "Do not scare them for sport. Frightened people send money. Clear people hang up.")

    # 13 Numbers
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(1.1), [
        ("The headlines — and the line they skip", 30, True, NAVY),
        ("FTC and FBI are different report piles. Name the source. Do not add them.", 16, False, SLATE),
    ])
    rect(s, Inches(0.6), Inches(1.55), Inches(6.0), Inches(4.7), NAVY)
    textbox(s, Inches(0.85), Inches(1.75), Inches(5.5), Inches(4.3), [
        ("FTC  ·  adults 60+", 16, True, ORANGE),
        ("$2.4 billion reported lost in 2024", 22, True, WHITE),
        ("About 4× the 2020 figure", 18, False, RGBColor(0xC8, 0xDF, 0xF0)),
        ("About 3 in 4 who reported a scam lost no money. They spotted it.", 20, True, GREEN),
    ])
    rect(s, Inches(6.85), Inches(1.55), Inches(5.85), Inches(4.7), INDIGO)
    textbox(s, Inches(7.1), Inches(1.75), Inches(5.4), Inches(4.3), [
        ("FBI IC3  ·  over 60", 16, True, CYAN),
        ("Nearly $5 billion in 2024", 22, True, WHITE),
        ("Most complaints of any age group", 18, False, RGBColor(0xC8, 0xDF, 0xF0)),
        ("Investment schemes took the most money. Fake tech support took almost a billion.", 18, False, WHITE),
    ])
    footer(s, 13, total, light=True)
    notes(s, "The dollars exploded because a smaller number of scams take life-changing amounts. They are not the punchline.")

    # 14 Spot a fake
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.8), [("Skill 1  ·  Spot a fake", 34, True, NAVY)])
    checks = [
        ("Who started this?", "If you did not start the conversation, slow down."),
        ("Does the address look odd?", "Look at the real email, not the logo."),
        ("Are they sending you somewhere?", "Link, QR code, or a number they provided."),
        ("Would they already know your name?", "'Dear Customer' from your bank is a tell."),
    ]
    for i, (t, b) in enumerate(checks):
        x = Inches(0.6) + (i % 2) * Inches(6.3)
        y = Inches(1.35) + (i // 2) * Inches(2.5)
        rect(s, x, y, Inches(6.05), Inches(2.3), WHITE)
        rect(s, x, y, Inches(0.16), Inches(2.3), ORANGE)
        textbox(s, x + Inches(0.4), y + Inches(0.3), Inches(5.4), Inches(1.7), [(t, 22, True, NAVY), (b, 18, False, SLATE)])
    footer(s, 14, total, light=True)
    notes(s, "AI cleaned up the spelling. 'No typos' is no longer a safety test.")

    # 15 Pressure
    s = new_dark(prs)
    bar(s)
    textbox(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7), [("Skill 2  ·  Name the pressure", 32, True, WHITE)])
    levers = [
        ("Hurry", "Twenty minutes or you are arrested."),
        ("Fear", "Warrant. Virus. Grandchild in jail."),
        ("Secrecy", "Don't tell your spouse or the teller."),
        ("Authority", "Badge. Caller ID. Case number. Costume."),
        ("Protect it", "Never move money to protect it."),
    ]
    for i, (t, b) in enumerate(levers):
        x = Inches(0.4) + i * Inches(2.55)
        rect(s, x, Inches(1.4), Inches(2.45), Inches(4.8), INDIGO)
        textbox(s, x + Inches(0.12), Inches(1.65), Inches(2.2), Inches(1.2), [(t, 20, True, ORANGE)], align=PP_ALIGN.CENTER)
        textbox(s, x + Inches(0.14), Inches(3.0), Inches(2.16), Inches(2.8), [(b, 16, False, WHITE)], align=PP_ALIGN.CENTER)
    footer(s, 15, total)
    notes(s, "Gift cards, wire, crypto ATM, courier, remote access = stop. The FTC is tired of 'move money to protect it.'")

    # 16 Hang up
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(1.8), Inches(12), Inches(3.2), [
        ("Hang up.", 44, True, NAVY),
        ("Look up.", 44, True, NAVY),
        ("Call back.", 44, True, ORANGE),
    ], align=PP_ALIGN.CENTER)
    textbox(s, Inches(1.2), Inches(5.5), Inches(10.9), Inches(0.8), [
        ("Use the number on the statement or the back of the card — not the number in the text.", 18, False, SLATE),
    ], align=PP_ALIGN.CENTER)
    footer(s, 16, total, light=True)
    notes(s, "The real FTC will never threaten you, tell you to move money to protect it, or demand gift cards.")

    # 17 Voice + code word
    s = new_dark(prs)
    bar(s, ORANGE)
    textbox(s, Inches(0.6), Inches(0.4), Inches(12), Inches(1.3), [
        ("A familiar voice is no longer proof", 32, True, WHITE),
        ("AARP: most adults 50+ already worry about cloned voices. Pick a family code word this week.", 18, False, RGBColor(0xC8, 0xDF, 0xF0)),
    ])
    rect(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(4.2), INDIGO)
    textbox(s, Inches(0.95), Inches(2.4), Inches(11.4), Inches(3.6), [
        ("Family code word", 20, True, ORANGE),
        ("Known only to your people. Never posted online.", 22, False, WHITE),
        ("Used when money or panic shows up.", 22, False, WHITE),
        ("If they don't know it — hang up. Then dial the number in your phone.", 22, True, WHITE),
    ])
    footer(s, 17, total)
    notes(s, "Most impostors still pose as companies and agencies, not family. Guard both.")

    # 18 Passwords + 2FA
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8), [("Skill 3  ·  Passwords and the second lock", 30, True, NAVY)])
    rect(s, Inches(0.6), Inches(1.25), Inches(6.0), Inches(5.1), NAVY)
    textbox(s, Inches(0.85), Inches(1.5), Inches(5.5), Inches(4.6), [
        ("NIST, in plain English", 18, True, ORANGE),
        ("Longer is better — aim for 15+ characters.", 18, False, WHITE),
        ("A phrase beats P@ssw0rd!", 18, False, WHITE),
        ("Email password ≠ bank password.", 18, False, WHITE),
        ("A notebook in a drawer is allowed.", 18, False, WHITE),
    ])
    rect(s, Inches(6.85), Inches(1.25), Inches(5.85), Inches(5.1), INDIGO)
    textbox(s, Inches(7.1), Inches(1.5), Inches(5.4), Inches(4.6), [
        ("2FA  ·  two-factor", 18, True, CYAN),
        ("Password + something you have.", 18, False, WHITE),
        ("Turn it on: email first, then the bank.", 18, False, WHITE),
        ("Settings → Security → 2-step / multifactor.", 18, False, WHITE),
        ("Never read that code to anyone.", 20, True, ORANGE),
    ])
    footer(s, 18, total, light=True)
    notes(s, "CISA: any second lock is better than none. Anyone who asks you to read the code is a scammer — FTC.")

    # 19 If money moved
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7), [("If money already moved", 32, True, NAVY)])
    steps = [
        "1   Stop sending.",
        "2   Call the bank — number on the card.",
        "3   Say wire / gift card / crypto ATM if that is what you used.",
        "4   Change the email password. Turn on 2FA.",
        "5   Report:  reportfraud.ftc.gov   and   ic3.gov",
        "6   Tell one trusted person. Shame helps them.",
    ]
    textbox(s, Inches(0.8), Inches(1.2), Inches(11.8), Inches(5.4), [(t, 24, False, SLATE) for t in steps])
    footer(s, 19, total, light=True)
    notes(s, "This is not a moral failure. Recovery is not a promise. Silence is a guarantee of no chance.")

    # 20 Fridge 3
    s = new_dark(prs)
    bar(s, ORANGE)
    textbox(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4), [("FRIDGE LINE", 14, True, ORANGE)])
    textbox(s, Inches(0.6), Inches(2.0), Inches(12), Inches(3.4), [
        ("Unexpected  +  hurry", 36, True, WHITE),
        ("+  money or a code", 36, True, WHITE),
        ("=  hang up and verify", 40, True, ORANGE),
    ], align=PP_ALIGN.CENTER)
    footer(s, 20, total)
    notes(s, "Hold this. Photograph it. This is the whole lesson on one line.")

    # 21 Do this today
    s = new_light(prs)
    textbox(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7), [("Do this today", 34, True, NAVY)])
    todos = [
        ("Lesson 1", "Write one real question on a sticky note."),
        ("Lesson 2", "Ask it in one tool. Screenshot the useful answer."),
        ("Lesson 3", "Turn on 2FA on email. Set a family code word."),
        ("Always", "Hang up. Look up. Call back."),
    ]
    for i, (k, v) in enumerate(todos):
        y = Inches(1.25) + i * Inches(1.25)
        rect(s, Inches(0.6), y, Inches(12.1), Inches(1.12), WHITE)
        rect(s, Inches(0.6), y, Inches(0.16), Inches(1.12), CYAN if i % 2 else ORANGE)
        textbox(s, Inches(1.0), y + Inches(0.12), Inches(2.4), Inches(0.85), [(k, 18, True, ORANGE)])
        textbox(s, Inches(3.5), y + Inches(0.28), Inches(8.8), Inches(0.65), [(v, 22, False, NAVY)])
    footer(s, 21, total, light=True)

    # 22 Close
    s = new_dark(prs)
    rect(s, 0, 0, Inches(0.22), H, ORANGE)
    textbox(s, Inches(0.7), Inches(1.6), Inches(12), Inches(3.6), [
        ("Digital independence is not a privilege.", 28, False, RGBColor(0xC8, 0xDF, 0xF0)),
        ("It is your right.", 40, True, WHITE),
        ("You're not too late.", 28, True, ORANGE),
    ])
    textbox(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.2), [
        ("learnmoretechnologies.com/start-free-lesson/", 20, True, CYAN),
        ("hello@learnmoretechnologies.com  ·  reportfraud.ftc.gov  ·  ic3.gov", 16, False, RGBColor(0x9A, 0xB0, 0xC0)),
    ])
    footer(s, 22, total)
    notes(s, "No invented urgency. We do not sell like the people we just taught them to refuse.")

    prs.save(OUT)
    print(f"Wrote {OUT} ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    build()
